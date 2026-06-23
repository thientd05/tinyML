#!/usr/bin/env python3
"""Chụp ảnh từng slide trong `slides.html` rồi nhét full-bleed vào `slides.pptx`.

Khác với `make_slides.py` (vẽ lại bằng shape thuần), script này giữ nguyên 100%
giao diện HTML bằng cách render qua Chrome headless và chụp từng phần tử `.slide`.

    ./env/bin/python tools/html_to_pptx.py

Phụ thuộc: playwright (dùng google-chrome hệ thống) + python-pptx + Pillow.
"""
from __future__ import annotations

import os
import tempfile

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
HTML = os.path.join(ROOT, "slides.html")
OUT = os.path.join(ROOT, "slides.pptx")

SLIDE_W_PX = 1280          # viewport 16:9 -> mỗi .slide cao 100vh = 720
SLIDE_H_PX = 720
SCALE = 2                  # device pixel ratio -> ảnh nét gấp đôi

# Ép hiện toàn bộ nội dung: tắt animation/transition + bỏ reveal opacity-0.
# Nếu không, slide chưa "active" sẽ chụp ra mờ hoặc trống.
FORCE_VISIBLE = """
*{animation:none!important;transition:none!important}
.rise{opacity:1!important;transform:none!important}
.trace.anim,.ecgline{stroke-dashoffset:0!important}
.counter,.fsbtn,.toprog{display:none!important}
"""


def capture(png_dir: str) -> list[str]:
    paths: list[str] = []
    with sync_playwright() as p:
        browser = p.chromium.launch(channel="chrome", headless=True)
        page = browser.new_page(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
            device_scale_factor=SCALE,
        )
        page.goto(f"file://{HTML}", wait_until="networkidle")
        page.add_style_tag(content=FORCE_VISIBLE)
        # đánh dấu mọi slide là active (phòng các kiểu reveal dựa .is-active)
        page.evaluate(
            "document.querySelectorAll('.slide').forEach(s=>s.classList.add('is-active'))"
        )
        page.evaluate("document.fonts.ready")
        page.wait_for_timeout(600)  # cho font ổn định

        slides = page.query_selector_all(".slide")
        print(f"  tìm thấy {len(slides)} slide")
        for i, el in enumerate(slides, 1):
            el.scroll_into_view_if_needed()
            page.wait_for_timeout(150)
            out = os.path.join(png_dir, f"slide_{i:02d}.png")
            el.screenshot(path=out)
            paths.append(out)
        browser.close()
    return paths


def build_pptx(pngs: list[str]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        with Image.open(png) as im:
            w, h = im.size
        ar_img = w / h
        ar_page = 13.333 / 7.5
        if ar_img >= ar_page:                 # ảnh rộng hơn -> khít chiều ngang
            pw = Inches(13.333)
            ph = Inches(13.333 / ar_img)
        else:                                 # ảnh cao hơn -> khít chiều dọc
            ph = Inches(7.5)
            pw = Inches(7.5 * ar_img)
        left = (Inches(13.333) - pw) // 2
        top = (Inches(7.5) - ph) // 2
        slide.shapes.add_picture(png, left, top, pw, ph)

    prs.save(OUT)


def main() -> None:
    with tempfile.TemporaryDirectory() as d:
        print("Chụp slide từ HTML...")
        pngs = capture(d)
        print(f"Dựng {OUT} từ {len(pngs)} ảnh...")
        build_pptx(pngs)
    print(f"Xong -> {OUT}")


if __name__ == "__main__":
    main()
