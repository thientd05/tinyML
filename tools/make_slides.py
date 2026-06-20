#!/usr/bin/env python3
"""Sinh slide thuyết trình `slide.pptx` (20 trang, theme sáng, tiếng Việt).

Nội dung lấy từ báo cáo `results/BaoCao_TinyML_ECG.docx`, `results/summary.csv`,
`results/winners.json` và các biểu đồ PNG trong `results/`. Chạy:

    ./env/bin/python tools/make_slides.py

Script độc lập, không phụ thuộc pipeline huấn luyện (chỉ cần python-pptx + Pillow).
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RES = os.path.join(ROOT, "results")
OUT = os.path.join(ROOT, "slide.pptx")

# ----------------------------------------------------------------------------- theme
INK = RGBColor(0x1A, 0x23, 0x30)        # chữ chính
MUTED = RGBColor(0x5B, 0x6B, 0x7C)      # chữ phụ
BG = RGBColor(0xFF, 0xFF, 0xFF)         # nền trắng
PANEL = RGBColor(0xF4, 0xF7, 0xFA)      # nền panel nhạt
ACCENT = RGBColor(0x0E, 0x7C, 0x86)     # xanh y tế
ACCENT2 = RGBColor(0x13, 0xA4, 0xB0)    # teal sáng
BLUE = RGBColor(0x25, 0x63, 0xA8)       # xanh dương phụ
WARN = RGBColor(0xE0, 0x8A, 0x1E)       # cam cảnh báo
WIN = RGBColor(0x2E, 0x8B, 0x57)        # xanh lá winner
WIN_BG = RGBColor(0xDD, 0xF0, 0xE4)     # nền hàng winner
HEAD_BG = RGBColor(0x0E, 0x7C, 0x86)    # nền header bảng
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# ----------------------------------------------------------------------------- helpers
def _solid(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def _run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    return r


def new_slide(bg=BG):
    slide = prs.slides.add_slide(BLANK)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _solid(rect, bg)
    _page["n"] += 1
    return slide


def footer(slide):
    # số trang + nhãn dự án
    tb, tf = _box(slide, Inches(11.7), Inches(7.02), Inches(1.5), Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{_page['n']:02d} / 20", 11, MUTED)
    tb2, tf2 = _box(slide, Inches(0.55), Inches(7.02), Inches(6), Inches(0.35))
    _run(tf2.paragraphs[0], "TinyML ECG · ESP32", 10, MUTED)


def header(slide, title, kicker=None):
    # thanh accent mảnh
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.55),
                                 Inches(0.12), Inches(0.62))
    _solid(bar, ACCENT)
    tb, tf = _box(slide, Inches(0.85), Inches(0.42), Inches(11.9), Inches(0.95))
    if kicker:
        pk = tf.paragraphs[0]
        _run(pk, kicker.upper(), 12, ACCENT2, bold=True)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    _run(p, title, 28, INK, bold=True)
    # đường kẻ mảnh dưới header
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.5),
                                Inches(12.23), Pt(1.5))
    _solid(ln, RGBColor(0xE2, 0xE8, 0xEE))
    footer(slide)


def bullets(slide, items, left, top, width, height, size=18, gap=10):
    """items: list of (text, level, color, bold)."""
    tb, tf = _box(slide, left, top, width, height)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, it in enumerate(items):
        text = it[0]
        level = it[1] if len(it) > 1 else 0
        color = it[2] if len(it) > 2 else INK
        bold = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        bull = "●  " if level == 0 else "–  "
        _run(p, bull, size - 2, ACCENT if level == 0 else MUTED, bold=True)
        _run(p, text, size, color, bold=bold)
    return tb


def panel(slide, l, t, w, h, color=PANEL):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _solid(sh, color)
    try:
        sh.adjustments[0] = 0.04
    except Exception:
        pass
    return sh


def image_fit(slide, path, l, t, max_w, max_h, caption=None):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    w, h = max_w, Emu(int(max_w / ar))
    if h > max_h:
        h, w = max_h, Emu(int(max_h * ar))
    left = l + Emu(int((max_w - w) / 2))
    top = t + Emu(int((max_h - h) / 2))
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    if caption:
        cb, cf = _box(slide, l, t + max_h + Inches(0.02), max_w, Inches(0.4))
        p = cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, caption, 12, MUTED, italic=True)


def table(slide, headers, rows, l, t, w, col_w=None, highlight=None, font=12,
          head_font=12, row_h=0.42):
    nrows, ncols = len(rows) + 1, len(headers)
    h = Inches(row_h * nrows)
    gf = slide.shapes.add_table(nrows, ncols, l, t, w, h).table
    if col_w:
        total = sum(col_w)
        for j, cwj in enumerate(col_w):
            gf.columns[j].width = Emu(int(int(w) * cwj / total))
    # header
    for j, htext in enumerate(headers):
        c = gf.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = HEAD_BG
        c.margin_top = Pt(2)
        c.margin_bottom = Pt(2)
        c.margin_left = Pt(6)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        _run(p, htext, head_font, WHITE, bold=True)
    # body
    for i, row in enumerate(rows, start=1):
        hl = highlight is not None and (i - 1) == highlight
        for j, val in enumerate(row):
            c = gf.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WIN_BG if hl else (PANEL if i % 2 else WHITE)
            c.margin_top = Pt(1)
            c.margin_bottom = Pt(1)
            c.margin_left = Pt(6)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            _run(p, str(val), font, WIN if hl else INK, bold=hl and j == 0)
    return gf


def chip(slide, text, l, t, w, color=ACCENT, txt=WHITE, h=0.5, size=14):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(h))
    _solid(sh, color)
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, size, txt, bold=True)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


# =========================================================================== SLIDES
def s01_title():
    slide = new_slide(BG)
    # khối màu trang trí trái
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), EMU_H)
    _solid(band, ACCENT)

    # chỉ còn tiêu đề, căn giữa theo chiều dọc
    tb, tf = _box(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(2.2))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _run(p, "Phát hiện bất thường nhịp tim ECG\n", 40, INK, bold=True)
    p2 = tf.add_paragraph()
    _run(p2, "bằng TinyML trên ESP32", 40, ACCENT, bold=True)

    footer(slide)


def s02_context():
    slide = new_slide()
    header(slide, "Bối cảnh & Mục tiêu", "Vấn đề")
    bullets(slide, [
        ("Bệnh tim mạch là nguyên nhân tử vong hàng đầu", 0, INK, True),
        ("Cần giám sát ECG liên tục, ngay tại thiết bị đeo (tại biên)", 1),
        ("Gửi tín hiệu lên cloud = tốn pin, trễ, lệ thuộc mạng", 1),
        ("Bài toán: phát hiện nhịp BẤT THƯỜNG ngay trên vi điều khiển", 0, INK, True),
        ("Bỏ sót bất thường là lỗi đắt nhất → ưu tiên recall cao", 1, WARN, False),
    ], Inches(0.7), Inches(1.75), Inches(6.7), Inches(4.8), size=18, gap=12)

    panel(slide, Inches(7.7), Inches(1.85), Inches(5.1), Inches(4.5))
    tb, tf = _box(slide, Inches(8.0), Inches(2.05), Inches(4.5), Inches(0.5))
    _run(tf.paragraphs[0], "MỤC TIÊU", 16, ACCENT, bold=True)
    bullets(slide, [
        ("Chọn mô hình tốt nhất dưới 3 ràng buộc phần cứng:", 0, INK, True),
        ("Độ trễ  ≤ 100 ms / nhịp", 1, BLUE, True),
        ("Flash  ≤ 1 MB / mô hình", 1, BLUE, True),
        ("RAM làm việc  ≤ 64 KB", 1, BLUE, True),
        ("Đạt recall ≥ 0.95 (độ nhạy lâm sàng)", 0, INK, True),
        ("Tối đa precision tại điểm vận hành đó", 1),
    ], Inches(8.0), Inches(2.55), Inches(4.5), Inches(3.6), size=16, gap=9)


def s03_team():
    slide = new_slide()
    header(slide, "Thành viên & Phân công", "Nhóm thực hiện")
    members = [
        ("Nguyễn Minh Hiếu", "Random Forest", "RF", ACCENT),
        ("Hàn Kim Trường", "Support Vector Machine", "SVM", BLUE),
        ("Tạ Đình Thiện", "1D Convolutional NN", "CNN", ACCENT2),
        ("Nguyễn Huy Vũ", "LSTM Recurrent NN", "LSTM", WARN),
    ]
    x = Inches(0.7)
    for name, role, tag, col in members:
        panel(slide, x, Inches(2.1), Inches(2.85), Inches(3.4))
        chip(slide, tag, x + Inches(0.42), Inches(2.45), Inches(2.0), color=col, h=0.7, size=22)
        tb, tf = _box(slide, x + Inches(0.15), Inches(3.55), Inches(2.55), Inches(1.6))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, name, 17, INK, bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        _run(p2, role, 14, MUTED)
        x = Emu(int(x) + int(Inches(3.04)))
    tb, tf = _box(slide, Inches(0.7), Inches(5.9), Inches(12), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Chung: pipeline dữ liệu · phương pháp tuyển chọn · triển khai & benchmark ESP32",
         15, MUTED, italic=True)


def s04_ecg():
    slide = new_slide()
    header(slide, "Tín hiệu ECG cơ bản", "Nền tảng")
    bullets(slide, [
        ("ECG = điện tâm đồ, ghi hoạt động điện của tim", 0, INK, True),
        ("Mỗi nhịp có phức bộ P – QRS – T", 1),
        ("R-peak: đỉnh nhọn nhất, dùng để cắt & căn nhịp", 1, BLUE, True),
        ("Khoảng RR = thời gian giữa hai nhịp liên tiếp", 1),
        ("Bình thường vs Bất thường", 0, INK, True),
        ("Hình dạng QRS & nhịp lệch → loạn nhịp", 1),
        ("VEB (thất): QRS rộng, méo → DỄ phát hiện", 1, WIN, False),
        ("SVEB (nhĩ): gần giống nhịp thường → KHÓ", 1, WARN, False),
    ], Inches(0.7), Inches(1.75), Inches(6.2), Inches(5.0), size=17, gap=10)
    image_fit(slide, os.path.join(RES, "beat_morphology_gallery.png"),
              Inches(7.1), Inches(1.8), Inches(5.7), Inches(4.6),
              caption="Hình dạng nhịp thực tế từ MIT-BIH (N · VEB · SVEB · block)")


def s05_dataset1():
    slide = new_slide()
    header(slide, "Bộ dữ liệu — MIT-BIH Arrhythmia", "Dữ liệu (1/3)")
    bullets(slide, [
        ("Chuẩn vàng cho nghiên cứu loạn nhịp", 0, INK, True),
        ("48 bản ghi · 2 kênh · tần số 360 Hz", 1),
        ("Dùng lead MLII (chọn theo TÊN kênh)", 1, BLUE, True),
        ("Bài toán nhị phân: Bình thường (0) vs Bất thường (1)", 0, INK, True),
        ("Nhãn AAMI: N/L/R/e/j → 0 ; còn lại → 1", 1),
        ("Mất cân bằng lớp: chỉ ~11% là bất thường", 1, WARN, True),
    ], Inches(0.7), Inches(1.75), Inches(6.4), Inches(4.8), size=18, gap=12)

    panel(slide, Inches(7.5), Inches(1.85), Inches(5.3), Inches(4.6))
    tb, tf = _box(slide, Inches(7.8), Inches(2.05), Inches(4.7), Inches(0.5))
    _run(tf.paragraphs[0], "CHIA THEO BỆNH NHÂN (de Chazal)", 15, ACCENT, bold=True)
    bullets(slide, [
        ("DS1 — huấn luyện (22 bản ghi)", 0, INK, True),
        ("DS1-val — bản ghi 207/209/215", 0, BLUE, True),
        ("dùng để chỉnh ngưỡng (patient-pure)", 1),
        ("DS2 — kiểm thử, bệnh nhân HOÀN TOÀN mới", 0, INK, True),
        ("Loại 4 bản ghi máy tạo nhịp (102/104/107/217)", 0, MUTED, False),
        ("Không trộn DS1+DS2 → tránh rò rỉ bệnh nhân", 1, WARN, True),
    ], Inches(7.8), Inches(2.6), Inches(4.7), Inches(3.6), size=15, gap=9)


def s06_cleaning():
    slide = new_slide()
    header(slide, "Làm sạch dữ liệu — KHÔNG dùng nhãn", "Dữ liệu (2/3)")
    table(slide,
          ["#", "Bước xử lý", "Mục đích"],
          [["1", "Chọn lead MLII theo tên", "Sửa bản ghi 114 (kênh 0 là V5)"],
           ["2", "RR đo giữa các nhịp THẬT", "Bỏ ~4.2% RR nhiễu bởi mốc phi-nhịp"],
           ["3", "Căn lại R-peak (±15 mẫu)", "Cửa sổ bám đúng đỉnh sóng"],
           ["4", "Lọc artifact (flatline/clipping)", "Bỏ nhịp hỏng, KHÔNG theo biên độ"]],
          Inches(0.7), Inches(1.8), Inches(8.2), col_w=[0.6, 4, 5.4],
          font=15, head_font=14, row_h=0.78)

    panel(slide, Inches(9.2), Inches(1.85), Inches(3.6), Inches(4.5), PANEL)
    tb, tf = _box(slide, Inches(9.45), Inches(2.05), Inches(3.1), Inches(4.2))
    p = tf.paragraphs[0]
    _run(p, "VÌ SAO QUAN TRỌNG", 14, ACCENT, bold=True)
    for t, c, b in [("Mọi bước đều label-free", BLUE, True),
                    ("→ không rò rỉ thông tin nhãn", MUTED, False),
                    ("Lọc theo VẬT LÝ, không theo biên độ", INK, True),
                    ("→ tránh thiên lệch lớp", MUTED, False),
                    ("Áp dụng đồng nhất mọi split", INK, True),
                    ("Kết quả: precision +0.017…0.023", WIN, True)]:
        pp = tf.add_paragraph()
        pp.space_before = Pt(8)
        _run(pp, t, 13, c, bold=b)


def s07_dataset3():
    slide = new_slide()
    header(slide, "Quy mô dữ liệu & Đặc trưng đầu vào", "Dữ liệu (3/3)")
    # stat row
    stats = [("100,697", "nhịp sau làm sạch", ACCENT),
             ("~42", "nhịp bị loại (~0.04%)", BLUE),
             ("~11%", "tỉ lệ bất thường", WARN),
             ("+0.02", "precision nhờ làm sạch", WIN)]
    x = Inches(0.7)
    for big, small, col in stats:
        panel(slide, x, Inches(1.8), Inches(2.9), Inches(1.3))
        cb, cf = _box(slide, x, Inches(1.95), Inches(2.9), Inches(1.05))
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, big, 24, col, bold=True)
        p2 = cf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, small, 12, MUTED)
        x = Emu(int(x) + int(Inches(3.05)))

    tb, tf = _box(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.5))
    _run(tf.paragraphs[0], "Mỗi họ mô hình nhận một dạng đầu vào riêng:", 17, INK, bold=True)
    table(slide,
          ["Đầu vào", "Mô tả", "Dùng cho"],
          [["21 đặc trưng tay", "thống kê thời gian · wavelet db4 · FFT · RR", "RF · SVM"],
           ["Raw 200×1", "nhịp thô ±100 mẫu quanh R-peak (~555 ms)", "CNN"],
           ["Raw 100×1", "nhịp thô giảm mẫu ×2", "LSTM"]],
          Inches(0.7), Inches(4.0), Inches(12.1), col_w=[2.4, 6.5, 2.2],
          font=15, head_font=14, row_h=0.62)


def family_intro(slide_title, kicker, idea_bullets, headers, rows, highlight,
                 note=None):
    slide = new_slide()
    header(slide, slide_title, kicker)
    bullets(slide, idea_bullets, Inches(0.7), Inches(1.75), Inches(4.9), Inches(5.0),
            size=17, gap=11)
    table(slide, headers, rows, Inches(5.9), Inches(1.85), Inches(6.9),
          col_w=[2.2, 1.3, 1.3, 1.6], font=12, head_font=12, row_h=0.5,
          highlight=highlight)
    if note:
        tb, tf = _box(slide, Inches(5.9), Inches(1.85) + Inches(0.5 * (len(rows) + 1)) + Inches(0.2),
                      Inches(6.9), Inches(0.9))
        _run(tf.paragraphs[0], note, 13, MUTED, italic=True)
    return slide


def family_result(slide_title, kicker, win_label, metrics, takeaways, badge_color):
    slide = new_slide()
    header(slide, slide_title, kicker)
    # khối winner
    panel(slide, Inches(0.7), Inches(1.85), Inches(5.4), Inches(4.5))
    chip(slide, "WINNER:  " + win_label, Inches(1.0), Inches(2.1), Inches(4.8),
         color=badge_color, h=0.6, size=18)
    # metric grid 2x2... dùng list
    tb, tf = _box(slide, Inches(1.05), Inches(2.95), Inches(4.8), Inches(3.2))
    for i, (k, v, col) in enumerate(metrics):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        _run(p, k + "   ", 15, MUTED)
        _run(p, v, 17, col, bold=True)
    # takeaways phải
    panel(slide, Inches(6.4), Inches(1.85), Inches(6.4), Inches(4.5), PANEL)
    tb, tf = _box(slide, Inches(6.75), Inches(2.1), Inches(5.8), Inches(0.5))
    _run(tf.paragraphs[0], "NHẬN XÉT", 15, ACCENT, bold=True)
    bullets(slide, takeaways, Inches(6.75), Inches(2.6), Inches(5.8), Inches(3.6),
            size=16, gap=12)
    return slide


def s08_09_rf():
    family_intro(
        "Random Forest — Thiết kế", "Họ mô hình · RF",
        [("Tập hợp nhiều cây quyết định", 0, INK, True),
         ("Đầu vào: 21 đặc trưng tay", 1),
         ("Bỏ phiếu đa số → xác suất bất thường", 1),
         ("Knob năng lực: tổng số nút cây", 0, INK, True),
         ("(n_estimators × max_depth)", 1, MUTED, False),
         ("Sweep 6 điểm: n10_d3 → n80_d12", 1, BLUE, True),
         ("Ưu: nhanh, rẻ, dễ giải thích", 0, WIN, False)],
        ["Cấu hình", "MACs", "Params", "Flash KB"],
        [["n10_d3", "30", "150", "2.9"],
         ["n10_d5  ★", "50", "612", "12.0"],
         ["n20_d6", "120", "2,170", "42.4"],
         ["n30_d8", "240", "8,228", "160.7"],
         ["n50_d10", "500", "23,484", "458.7"],
         ["n80_d12", "960", "51,468", "1,005"]],
        highlight=1,
        note="★ = cấu hình thắng trong họ.  n80_d12 ~ vượt ngưỡng flash 1 MB.")
    family_result(
        "Random Forest — Kết quả", "Họ mô hình · RF", "rf_n10_d5",
        [("Precision @ recall 0.95", "0.127 ± 0.01", INK),
         ("PR-AUC", "0.742", INK),
         ("Độ trễ / Flash", "1.0 ms · 12 KB", BLUE),
         ("FPR @ điểm vận hành", "0.835", WARN),
         ("Recall khi triển khai (DS2)", "0.78  ⚠", WARN)],
        [("Mô hình nhỏ đã đạt chất lượng tốt nhất", 0, INK, True),
         ("To hơn KHÔNG cải thiện precision", 1),
         ("PR-AUC cao nhưng…", 0, INK, True),
         ("ngưỡng chỉnh trên DS1 tụt mạnh trên DS2", 1, WARN, True),
         ("→ vấn đề chuyển ngưỡng (threshold transfer)", 1, WARN, False)],
        ACCENT)


def s10_11_svm():
    family_intro(
        "SVM — Thiết kế", "Họ mô hình · SVM",
        [("Tìm siêu phẳng phân tách tối ưu", 0, INK, True),
         ("Đầu vào: 21 đặc trưng tay (như RF)", 1),
         ("Linear vs RBF kernel", 1),
         ("Knob năng lực: số support vector", 0, INK, True),
         ("linear = 1 vector trọng số (21 số)", 1, BLUE, True),
         ("Sweep: linear, rbf1k … rbf10k", 1),
         ("Class-weight cân bằng cho lớp hiếm", 0, INK, True)],
        ["Cấu hình", "MACs", "Params", "Flash KB"],
        [["linear  ★", "21", "21", "0.16"],
         ["rbf1k", "3,612", "3,612", "14.9"],
         ["rbf2k", "6,069", "6,069", "25.0"],
         ["rbf4k", "9,492", "9,492", "39.0"],
         ["rbf7k", "18,123", "18,123", "74.3"],
         ["rbf10k", "22,911", "22,911", "93.9"]],
        highlight=0,
        note="★ linear: nhỏ nhất, xác định (deterministic), tái lập hoàn toàn.")
    family_result(
        "SVM — Kết quả", "Họ mô hình · SVM", "svm_linear  ⭐ QUÁN QUÂN",
        [("Precision @ recall 0.95", "0.193 ± 0.00", WIN),
         ("PR-AUC", "0.730", INK),
         ("Độ trễ / Flash", "1.0 ms · 0.2 KB", BLUE),
         ("FPR @ điểm vận hành", "0.49  (thấp nhất)", WIN),
         ("Recall khi triển khai (DS2)", "0.92  ⚠", WARN)],
        [("Precision TRUNG BÌNH cao nhất toàn cục", 0, WIN, True),
         ("std = 0 → hoàn toàn tái lập", 1),
         ("Nhỏ & đơn giản nhất: 21 số, 0.2 KB", 0, INK, True),
         ("FPR thấp nhất trong các winner", 0, BLUE, True),
         ("Caveat: recall DS2 0.92 (chuyển ngưỡng)", 1, WARN, False)],
        BLUE)


def s12_13_cnn():
    family_intro(
        "1D-CNN — Thiết kế", "Họ mô hình · CNN",
        [("Học đặc trưng trực tiếp từ nhịp thô", 0, INK, True),
         ("Đầu vào: raw 200×1 (không cần feature tay)", 1, BLUE, True),
         ("Khối Conv1D + BatchNorm + MaxPool", 1),
         ("Knob năng lực: số MACs / nhịp", 0, INK, True),
         ("Sweep: c4 → c16-32-64-64", 1),
         ("GPU train, mixed precision", 1, MUTED, False),
         ("Chuyển ngưỡng DS1→DS2 ổn định", 0, WIN, False)],
        ["Cấu hình", "MACs", "Params", "Flash KB"],
        [["c4", "4,008", "34", "0.1"],
         ["c8", "8,016", "66", "0.3"],
         ["c8-16  ★", "84,832", "13,554", "52.9"],
         ["c16-16", "156,832", "14,242", "55.6"],
         ["c16-32-32", "553,664", "33,538", "131"],
         ["c16-32-64-64", "1.35 M", "82,882", "324"]],
        highlight=2,
        note="★ c8-16.  Biến thể lớn nhất: 239 ms → VI PHẠM ngưỡng 100 ms ❌.")
    family_result(
        "1D-CNN — Kết quả", "Họ mô hình · CNN", "cnn_c8-16",
        [("Precision @ recall 0.95", "0.172 ± 0.05", WARN),
         ("PR-AUC", "0.635", INK),
         ("Độ trễ / Flash", "19.8 ms · 52.9 KB", BLUE),
         ("Recall khi triển khai (DS2)", "0.965  ✓", WIN),
         ("FPR @ điểm vận hành", "0.632", INK)],
        [("Chuyển ngưỡng DS1→DS2 RẤT tốt (recall 0.965)", 0, WIN, True),
         ("Nhưng phương sai theo seed LỚN (±0.05)", 0, WARN, True),
         ("1 lần chạy may mắn ≠ chất lượng thật", 1),
         ("→ phải so sánh bằng trung bình 5 seed", 1, MUTED, False),
         ("Biến thể lớn vi phạm độ trễ & không khả thi", 0, WARN, False)],
        ACCENT2)


def s14_15_lstm():
    family_intro(
        "LSTM — Thiết kế", "Họ mô hình · LSTM",
        [("Mạng hồi tiếp, mô hình hóa chuỗi thời gian", 0, INK, True),
         ("Đầu vào: raw 100×1 (giảm mẫu ×2)", 1, BLUE, True),
         ("Cổng nhớ nắm bắt phụ thuộc dài", 1),
         ("Knob năng lực: hidden size", 0, INK, True),
         ("Sweep: h4 … h32, h32x2 (2 lớp)", 1),
         ("Đủ nhỏ để chạy fp32", 1, MUTED, False),
         ("Chuyển ngưỡng tốt nhất trong 4 họ", 0, WIN, False)],
        ["Cấu hình", "MACs", "Params", "Flash KB"],
        [["h4", "8,008", "122", "0.5"],
         ["h8  ★", "28,816", "370", "1.4"],
         ["h16", "108,832", "1,250", "4.9"],
         ["h24", "240,048", "2,642", "10.3"],
         ["h32", "422,464", "4,546", "17.8"],
         ["h32x2", "1.24 M", "12,994", "50.8"]],
        highlight=1,
        note="★ h8.  h32x2: 262 ms → VI PHẠM ngưỡng 100 ms ❌.")
    family_result(
        "LSTM — Kết quả", "Họ mô hình · LSTM", "lstm_h8",
        [("Precision @ recall 0.95", "0.152 ± 0.01", INK),
         ("PR-AUC", "0.476  (thấp nhất)", WARN),
         ("Độ trễ / Flash", "31.1 ms · 1.4 KB", BLUE),
         ("Recall khi triển khai (DS2)", "0.987  ✓", WIN),
         ("FPR @ điểm vận hành", "0.657", INK)],
        [("Chuyển ngưỡng tốt nhất (recall 0.987)", 0, WIN, True),
         ("Flash rất nhỏ (1.4 KB)", 0, BLUE, False),
         ("Nhưng PR-AUC thấp nhất → tách lớp yếu hơn", 0, WARN, True),
         ("Độ trễ cao hơn CNN cùng chất lượng", 1),
         ("Mô hình nhỏ (h8) vẫn là tốt nhất trong họ", 0, INK, False)],
        WARN)


def s16_summary_table():
    slide = new_slide()
    header(slide, "So sánh tổng hợp — Winner của 4 họ", "Kết quả (1/5)")
    table(slide,
          ["Họ", "Winner", "Precision\n@rec 0.95", "FPR", "PR-AUC",
           "Trễ (ms)", "Flash KB", "Recall DS2", "Khả thi"],
          [["RF", "n10_d5", "0.127", "0.83", "0.742", "1.0", "12.0", "0.78 ⚠", "✓"],
           ["SVM", "linear", "0.193", "0.49", "0.730", "1.0", "0.2", "0.92 ⚠", "✓"],
           ["CNN", "c8-16", "0.172", "0.63", "0.635", "19.8", "52.9", "0.965 ✓", "✓"],
           ["LSTM", "h8", "0.152", "0.66", "0.476", "31.1", "1.4", "0.987 ✓", "✓"]],
          Inches(0.7), Inches(2.0), Inches(12.1),
          col_w=[0.9, 1.5, 1.6, 1.0, 1.2, 1.2, 1.2, 1.5, 1.0],
          font=14, head_font=12.5, row_h=0.62, highlight=1)
    panel(slide, Inches(0.7), Inches(5.3), Inches(12.1), Inches(1.3), PANEL)
    tb, tf = _box(slide, Inches(1.0), Inches(5.5), Inches(11.6), Inches(1.0))
    p = tf.paragraphs[0]
    _run(p, "Quán quân toàn cục:  ", 18, INK, bold=True)
    _run(p, "SVM linear", 18, WIN, bold=True)
    _run(p, "  — precision trung bình cao nhất, nhanh nhất, nhỏ nhất, xác định.", 18, INK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    _run(p2, "⚠ RF/SVM tụt recall trên DS2 (chuyển ngưỡng);  CNN/LSTM chuyển ngưỡng tốt nhưng đắt hơn.",
         14, MUTED, italic=True)


def s17_quality_cost():
    slide = new_slide()
    header(slide, "Chất lượng vs Chi phí — “To hơn KHÔNG tốt hơn”", "Kết quả (2/5)")
    image_fit(slide, os.path.join(RES, "within_family_quality_vs_cost.png"),
              Inches(0.7), Inches(1.75), Inches(8.7), Inches(4.9),
              caption="Precision & PR-AUC theo MACs (log) — vòng tròn = winner mỗi họ")
    panel(slide, Inches(9.6), Inches(1.85), Inches(3.2), Inches(4.5), PANEL)
    bullets(slide, [
        ("Chất lượng ĐẠT ĐỈNH ở mô hình nhỏ", 0, INK, True),
        ("Tăng năng lực → bão hòa / giảm", 1, WARN, False),
        ("Mô hình lớn = đắt mà không lợi", 0, INK, True),
        ("Chọn nhỏ: rẻ, nhanh, khả thi", 0, WIN, True),
    ], Inches(9.8), Inches(2.1), Inches(2.85), Inches(4.0), size=15, gap=14)


def s18_pareto():
    slide = new_slide()
    header(slide, "Mặt Pareto đa họ — chọn winner toàn cục", "Kết quả (3/5)")
    image_fit(slide, os.path.join(RES, "cross_family_pareto.png"),
              Inches(0.7), Inches(1.75), Inches(8.7), Inches(4.9),
              caption="Precision vs độ trễ (log) — vạch đứng = ngưỡng 100 ms")
    panel(slide, Inches(9.6), Inches(1.85), Inches(3.2), Inches(4.5), PANEL)
    bullets(slide, [
        ("Vùng khả thi = bên trái 100 ms", 0, INK, True),
        ("SVM linear ở góc TỐT NHẤT", 0, WIN, True),
        ("cao nhất & trái nhất", 1),
        ("CNN/LSTM lớn rơi ra ngoài ngưỡng", 0, WARN, False),
        ("Đơn giản nhất thắng theo trung bình", 0, INK, True),
    ], Inches(9.8), Inches(2.1), Inches(2.85), Inches(4.0), size=15, gap=13)


def s19_why_low():
    slide = new_slide()
    header(slide, "Vì sao precision thấp (0.13 – 0.19)?", "Kết quả (4/5)")
    bullets(slide, [
        ("Bài toán KHÓ về bản chất — không phải lỗi mô hình", 0, INK, True),
        ("Mất cân bằng lớp ~11% → baseline precision ~0.11", 1),
        ("Dịch phân bố giữa các bệnh nhân (DS1→DS2)", 1),
        ("SVEB gần giống nhịp thường → khó tách", 1, WARN, True),
        ("Mâu thuẫn nhãn AAMI: L/R rộng nhưng gán Normal", 1, WARN, False),
        ("Hướng khắc phục", 0, WIN, True),
        ("Debounce N nhịp liên tiếp · ngưỡng theo bệnh nhân", 1),
        ("Thêm lead V1 · tầng tin cậy đẩy ca khó lên cloud", 1),
    ], Inches(0.7), Inches(1.75), Inches(6.3), Inches(5.0), size=16, gap=10)
    image_fit(slide, os.path.join(RES, "beat_easy_vs_hard.png"),
              Inches(7.2), Inches(1.8), Inches(5.6), Inches(4.6),
              caption="Normal vs VEB (dễ) vs SVEB (chồng lấn, khó)")


def s20_conclusion():
    slide = new_slide()
    header(slide, "Kết luận & Triển khai ESP32", "Kết quả (5/5)")
    # stat row deployment
    stats = [("1.0000", "parity 24/24 mô hình", WIN),
             ("~5.5 KB", "RAM / 64 KB budget", BLUE),
             ("937 µs", "trích đặc trưng / nhịp", ACCENT),
             ("Compute", "là nút thắt, không phải RAM", WARN)]
    x = Inches(0.7)
    for big, small, col in stats:
        panel(slide, x, Inches(1.8), Inches(2.9), Inches(1.3))
        cb, cf = _box(slide, x, Inches(1.95), Inches(2.9), Inches(1.05))
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, big, 22, col, bold=True)
        p2 = cf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, small, 12, MUTED)
        x = Emu(int(x) + int(Inches(3.05)))

    panel(slide, Inches(0.7), Inches(3.35), Inches(7.0), Inches(3.0))
    tb, tf = _box(slide, Inches(1.0), Inches(3.55), Inches(6.5), Inches(0.5))
    _run(tf.paragraphs[0], "CHỐT HẠ", 15, ACCENT, bold=True)
    bullets(slide, [
        ("Winner: SVM linear", 0, WIN, True),
        ("precision 0.193 · 1 ms · 0.2 KB · tái lập", 1),
        ("Phương pháp: sweep năng lực + điểm vận hành", 0, INK, True),
        ("+ cổng khả thi phần cứng + trung bình 5 seed", 1),
        ("Cả pipeline PC ↔ ESP32 khớp bit-for-bit", 0, INK, True),
    ], Inches(1.0), Inches(4.05), Inches(6.5), Inches(2.2), size=15, gap=10)

    panel(slide, Inches(7.95), Inches(3.35), Inches(4.85), Inches(3.0), PANEL)
    tb, tf = _box(slide, Inches(8.25), Inches(3.55), Inches(4.3), Inches(0.5))
    _run(tf.paragraphs[0], "HƯỚNG PHÁT TRIỂN", 15, BLUE, bold=True)
    bullets(slide, [
        ("Lượng tử hóa / cắt tỉa (hiện fp32)", 0, INK, False),
        ("Trích đặc trưng on-device vào RF/SVM", 0, INK, False),
        ("Phân loại đa lớp AAMI (5 lớp)", 0, INK, False),
        ("Thêm lead & dữ liệu để nâng precision", 0, INK, False),
    ], Inches(8.25), Inches(4.05), Inches(4.3), Inches(2.2), size=15, gap=12)
    # thank you
    tb, tf = _box(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Cảm ơn đã lắng nghe!", 16, ACCENT, bold=True)


def main():
    s01_title()
    s02_context()
    s03_team()
    s04_ecg()
    s05_dataset1()
    s06_cleaning()
    s07_dataset3()
    s08_09_rf()
    s10_11_svm()
    s12_13_cnn()
    s14_15_lstm()
    s16_summary_table()
    s17_quality_cost()
    s18_pareto()
    s19_why_low()
    s20_conclusion()
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
